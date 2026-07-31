"""Render a deck spec to .pptx, following the company deck theme.

This is the offline path: it needs no Google connection, so it works when the
Slides API is unavailable or unauthorised. Uploading the result with conversion
turns it into native Google Slides.

Layout is measured, not hardcoded. python-pptx cannot ask a font how wide a
string is, so `_est_lines` simulates word wrap to predict how tall a block will
be, and every renderer flows downward from a cursor. Fixed offsets would collide
the moment real text wrapped to a second line.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from services.decks.archetypes import BY_NAME
from services.decks.theme import DEFAULT_THEME, EMU_PER_INCH, DeckTheme

logger = logging.getLogger(__name__)

# Average glyph width as a fraction of font size, for Archivo/JetBrains Mono.
# Deliberately generous: over-estimating costs a little whitespace, while
# under-estimating spills text onto whatever sits below it.
_CHAR_W = 0.52
_CHAR_W_BOLD = 0.56
_LINE_H = 1.25   # natural line box as a multiple of font size
_SAFETY = 1.06   # headroom for glyphs wider than the average

CONTENT_TOP = 1.02     # first baseline below the confidential mark
CONTENT_BOTTOM = 6.58  # last usable line above the page number

GAP_XS = 0.08
GAP_S = 0.16
GAP_M = 0.30
GAP_L = 0.44


class DeckBuildError(Exception):
    """A spec the caller should fix, reported verbatim."""


def _in(v: float) -> Emu:
    return Emu(int(round(v * EMU_PER_INCH)))


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _est_lines(text: str, width_in: float, size: float, bold: bool) -> int:
    """Simulate word wrap to predict line count."""
    if not text:
        return 1
    char_w = size * (_CHAR_W_BOLD if bold else _CHAR_W)
    per_line = max(1, int(width_in * 72 / char_w))
    total = 0
    for block in text.split("\n"):
        words = block.split()
        if not words:
            total += 1
            continue
        lines, cur = 1, 0
        for word in words:
            need = len(word) + (1 if cur else 0)
            if cur and cur + need > per_line:
                lines += 1
                cur = len(word)
            else:
                cur += need
        total += lines
    return total


def _line_pts(size: float, spacing: float) -> float:
    """Line box in points. `spacing` multiplies the natural box, not the font size."""
    return size * _LINE_H * spacing


def _text_h(lines: int, size: float, spacing: float) -> float:
    return lines * _line_pts(size, spacing) * _SAFETY / 72


class _SlideWriter:
    """Thin wrapper over a python-pptx slide that speaks the theme's language.

    Every text method returns the height it consumed, in inches, so callers can
    stack the next element beneath it.
    """

    def __init__(self, slide, theme: DeckTheme):
        self.s = slide
        self.t = theme

    # ── primitives ────────────────────────────────────────────────────

    def rect(self, left, top, width, height, fill: str):
        sh = self.s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _in(left), _in(top), _in(width), _in(height)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def rule(self, left, top, width, colour: str) -> float:
        self.rect(left, top, width, self.t.grid.hairline_h, colour)
        return self.t.grid.hairline_h

    def text(
        self,
        left,
        top,
        width,
        runs: list[tuple[str, str]] | str,
        *,
        font: str | None = None,
        size: float = 12,
        bold: bool = False,
        colour: str | None = None,
        align=PP_ALIGN.LEFT,
        spacing: float = 1.0,
        max_lines: int | None = None,
        min_size: float = 7.0,
    ) -> float:
        """Draw text sized to its own content. Returns height in inches.

        runs is either a plain string or [(text, colour_hex), …] for two-tone.
        max_lines shrinks the type until it fits rather than letting it overrun.
        """
        if isinstance(runs, str):
            runs = [(runs, colour or self.t.palette.primary)]
        runs = [(t, c) for t, c in runs if t]
        if not runs:
            return 0.0

        flat = "".join(t for t, _ in runs)
        if max_lines:
            while size > min_size and _est_lines(flat, width, size, bold) > max_lines:
                size -= 0.5
        lines = _est_lines(flat, width, size, bold)
        height = _text_h(lines, size, spacing)

        box = self.s.shapes.add_textbox(_in(left), _in(top), _in(width), _in(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        para = tf.paragraphs[0]
        para.alignment = align
        # Absolute points, not a percentage: Google Slides reads spcPct against
        # the font size and PowerPoint against the line box, so a percentage
        # renders differently in each and collapses lines in one of them.
        para.line_spacing = Pt(_line_pts(size, spacing))
        for txt, col in runs:
            run = para.add_run()
            run.text = txt
            run.font.name = font or self.t.fonts.display
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = _rgb(col)
        return height

    # ── theme idioms ──────────────────────────────────────────────────

    def mono(self, left, top, width, txt, *, colour=None, size=None,
             align=PP_ALIGN.LEFT, max_lines=1) -> float:
        return self.text(
            left, top, width, str(txt),
            font=self.t.fonts.mono,
            size=size or self.t.sizes.eyebrow,
            colour=colour or self.t.palette.dim,
            align=align,
            spacing=1.0,
            max_lines=max_lines,
        )

    def chrome(self, page_no: int | None, classification: str = "") -> None:
        c, p, s = self.t.chrome, self.t.palette, self.t.sizes
        if classification:
            # The most sensitive tier is marked in the accent colour so it reads
            # as a warning rather than as furniture.
            colour = p.accent if "TRADE SECRET" in classification else p.dim
            self.mono(c.confidential_l, c.confidential_t, c.confidential_w,
                      classification, colour=colour)
        if page_no is not None:
            self.text(
                c.page_no_l, c.page_no_t, c.page_no_w, f"{page_no:02d}",
                font=self.t.fonts.mono, size=s.page_no, colour=p.dim,
            )

    def eyebrow(self, top: float, label: str) -> float:
        """Red tick plus an all-caps label. The tick sits on the text's midline."""
        if not label:
            return 0.0
        g, p = self.t.grid, self.t.palette
        height = self.mono(g.eyebrow_label_l, top, 4.0, str(label).upper(), colour=p.primary)
        self.rule(g.margin_l, top + height / 2, g.eyebrow_tick_w, p.accent)
        return height

    def headline(self, top, text, accent, *, size=None, width=None, max_lines=3) -> float:
        g, p, s = self.t.grid, self.t.palette, self.t.sizes
        runs = [(str(text), p.primary)]
        if accent:
            runs = [(str(text).rstrip() + " ", p.primary), (str(accent), p.accent)]
        return self.text(
            g.margin_l, top, width or g.content_w, runs,
            size=size or s.headline, bold=True, spacing=0.98, max_lines=max_lines,
        )

    def card(self, left, top, width, height, *, fill=None):
        return self.rect(left, top, width, height, fill or self.t.palette.card)

    # ── measuring, without drawing ─────────────────────────────────────

    def h(self, text, width, size, *, bold=False, spacing=1.0) -> float:
        return _text_h(_est_lines(str(text), width, size, bold), size, spacing)

    def h_mono(self, size=None) -> float:
        return _text_h(1, size or self.t.sizes.eyebrow, 1.0)


def _header(w: _SlideWriter, d: dict[str, Any], *, top: float = CONTENT_TOP,
            size: float | None = None, width: float | None = None,
            caption_colour: str | None = None) -> float:
    """Eyebrow, headline and optional caption. Returns the y below them."""
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = top
    eb = w.eyebrow(y, str(d.get("eyebrow", "")))
    if eb:
        y += eb + GAP_S
    y += w.headline(y, str(d.get("headline", "")), d.get("headline_accent"),
                    size=size, width=width)
    y += GAP_M
    if d.get("caption"):
        y += w.text(g.margin_l, y, width or g.rule_w, str(d["caption"]),
                    size=s.body_sm, colour=caption_colour or p.caption, max_lines=3)
        y += GAP_M
    return y


def _fits(w: _SlideWriter, name: str, y: float) -> None:
    if y > CONTENT_BOTTOM:
        logger.warning("Slide content overruns the canvas on %s (y=%.2fin)", name, y)


# ── archetype renderers ───────────────────────────────────────────────
# Each takes (writer, data) and flows content below the chrome.


def _cover(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    text_w = 7.60  # leaves the right third for the media panel
    y = 2.03
    if d.get("kicker"):
        h = w.mono(g.eyebrow_label_l, y, 4.0, str(d["kicker"]).upper(), colour=p.primary)
        w.rule(g.margin_l, y + h / 2, g.eyebrow_tick_w, p.accent)
        y += h + GAP_M
    y += w.text(g.margin_l, y, text_w, str(d.get("headline", "")),
                size=s.hero, bold=True, spacing=0.95, max_lines=2)
    y += GAP_M
    if d.get("subtitle"):
        runs = [(str(d["subtitle"]).rstrip() + " ", p.primary)]
        if d.get("subtitle_accent"):
            runs.append((str(d["subtitle_accent"]), p.accent))
        y += w.text(g.margin_l, y, text_w, runs, size=s.subtitle, max_lines=2)
        y += GAP_M
    if d.get("footnote"):
        y += w.text(g.margin_l, y, text_w, str(d["footnote"]),
                    font=w.t.fonts.mono, size=s.item, colour=p.caption, max_lines=2)
    _maybe_image(w, d, 8.42, 1.27, 3.31, 4.68)
    _fits(w, "cover", y)


def _section_break(w: _SlideWriter, d: dict[str, Any]) -> None:
    y = 2.03
    eb = w.eyebrow(y, str(d.get("eyebrow", "")))
    if eb:
        y += eb + GAP_M
    y += w.headline(y, str(d.get("headline", "")), d.get("headline_accent"),
                    size=w.t.sizes.headline_xl, max_lines=3)
    _fits(w, "section_break", y)


def _statement_media(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    media_l = 7.13
    text_w = media_l - g.margin_l - 0.38
    w.rect(media_l, 0.0, 5.70, 7.22, p.panel)
    _maybe_image(w, d, media_l, 0.0, 5.70, 7.22, label=str(d.get("media_label", "IMAGE")))

    y = 1.62
    eb = w.eyebrow(y, str(d.get("eyebrow", "")))
    if eb:
        y += eb + GAP_S
    y += w.headline(y, str(d.get("headline", "")), d.get("headline_accent"),
                    size=s.headline_sm, width=text_w, max_lines=6)
    y += GAP_L
    if d.get("caption"):
        y += w.text(g.margin_l, y, text_w, str(d["caption"]),
                    size=s.detail, colour=p.caption, max_lines=4)
    _fits(w, "statement_media", y)


def _bullets(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d)
    for item in [str(i) for i in (d.get("items") or [])]:
        y += w.rule(g.margin_l, y, g.rule_w, p.dim)
        y += GAP_S
        y += w.text(g.margin_l, y, g.rule_w, item, size=s.body, colour=p.body, max_lines=3)
        y += GAP_M
    _fits(w, "bullets", y)


def _points_with_metrics(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    left_w = 5.13
    y_left = CONTENT_TOP
    eb = w.eyebrow(y_left, str(d.get("eyebrow", "")))
    if eb:
        y_left += eb + GAP_S
    w.headline(y_left, str(d.get("headline", "")), d.get("headline_accent"),
               size=s.headline_xl, width=left_w, max_lines=4)

    y = CONTENT_TOP
    for i, item in enumerate([str(i) for i in (d.get("items") or [])][:4]):
        y += w.rule(g.right_l, y, g.right_w, p.primary)
        y += GAP_M
        w.text(g.right_l, y, 0.64, f"{i + 1:02d}", font=w.t.fonts.mono,
               size=s.index, colour=p.accent if i == 0 else p.dim)
        y += w.text(g.right_l + 0.81, y, g.right_w - 0.81, item,
                    size=s.body, colour=p.body, max_lines=3)
        y += GAP_M

    metrics = list(d.get("metrics") or [])[:2]
    if metrics:
        y += w.rule(g.right_l, y, g.right_w, p.primary)
        y += GAP_M
        col_w = g.right_w / len(metrics)
        tallest = 0.0
        for i, m in enumerate(metrics):
            x = g.right_l + i * col_w
            vh = w.text(x, y, col_w - 0.20, str(m.get("value", "")),
                        size=s.figure_lg, bold=True,
                        colour=p.accent if i == 0 else p.primary, max_lines=1)
            lh = w.mono(x, y + vh + GAP_XS, col_w - 0.20,
                        str(m.get("label", "")).upper(), colour=p.caption)
            tallest = max(tallest, vh + GAP_XS + lh)
        y += tallest
    _fits(w, "points_with_metrics", y)


def _numbered_cards(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    left_w = 5.42
    y_left = CONTENT_TOP
    eb = w.eyebrow(y_left, str(d.get("eyebrow", "")))
    if eb:
        y_left += eb + GAP_S
    y_left += w.headline(y_left, str(d.get("headline", "")), d.get("headline_accent"),
                         size=s.headline_xl, width=left_w, max_lines=4)
    if d.get("caption"):
        y_left += GAP_M
        w.text(g.margin_l, y_left, left_w - 0.40, str(d["caption"]),
               size=s.body_sm, colour=p.caption, max_lines=4)

    card_l, card_w = 6.46, 5.81
    inset = g.card_inset
    num_w = 0.56
    body_w = card_w - inset * 2 - num_w
    y = CONTENT_TOP
    for i, item in enumerate([str(i) for i in (d.get("items") or [])][:4]):
        card_h = max(0.94, w.h(item, body_w, s.body_sm) + inset * 2)
        w.card(card_l, y, card_w, card_h)
        w.text(card_l + inset, y + inset, num_w, f"{i + 1:02d}",
               font=w.t.fonts.mono, size=s.index,
               colour=p.accent if i == 0 else p.dim)
        w.text(card_l + inset + num_w, y + inset, body_w, item,
               size=s.body_sm, colour=p.primary)
        y += card_h + GAP_S
    _fits(w, "numbered_cards", y)


def _profile_cards(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d, size=s.headline_lg)

    cards = list(d.get("cards") or [])[:3]
    inner = g.col_3_w - g.card_inset * 2
    heights = [
        w.h_mono()
        + GAP_S + w.h(str(c.get("name", "")), inner, s.name, bold=True)
        + GAP_XS + w.h(str(c.get("detail", "")), inner, s.detail_sm)
        + g.card_inset * 2
        for c in cards
    ]
    card_h = max(heights) if heights else 0.0

    for i, c in enumerate(cards):
        x = g.cols_3[i]
        w.card(x, y, g.col_3_w, card_h)
        cy = y + g.card_inset
        cy += w.mono(x + g.card_inset, cy, inner, str(c.get("label", "")).upper(),
                     colour=p.accent)
        cy += GAP_S
        cy += w.text(x + g.card_inset, cy, inner, str(c.get("name", "")),
                     size=s.name, bold=True, max_lines=2)
        cy += GAP_XS
        w.text(x + g.card_inset, cy, inner, str(c.get("detail", "")),
               size=s.detail_sm, colour=p.caption, max_lines=5)
    y += card_h

    if d.get("footer_items"):
        y += GAP_L
        y += w.rule(g.margin_l, y, g.rule_w, p.primary)
        y += GAP_M
        label_w = 2.00
        if d.get("footer_label"):
            w.mono(g.margin_l, y, label_w, str(d["footer_label"]).upper(), colour=p.caption)
        items = list(d["footer_items"])[:4]
        x = g.margin_l + label_w + 0.10
        col_w = (g.rule_w - label_w - 0.10) / len(items)
        tallest = 0.0
        for item in items:
            if isinstance(item, dict):
                runs = [(str(item.get("name", "")) + " ", p.primary),
                        (str(item.get("detail", "")), p.dim)]
            else:
                runs = [(str(item), p.primary)]
            tallest = max(tallest, w.text(x, y, col_w - 0.12, runs,
                                          size=s.item, max_lines=2))
            x += col_w
        y += tallest
    _fits(w, "profile_cards", y)


def _tier_cards(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d, caption_colour=p.accent)

    cards = list(d.get("cards") or [])[:3]
    inner = g.col_3_w - g.card_inset * 2
    heights = []
    for c in cards:
        h = w.h_mono() + GAP_S
        h += w.h(str(c.get("value", "")), inner, s.figure_lg, bold=True) + GAP_M
        h += g.hairline_h + GAP_S
        for item in list(c.get("items") or [])[:5]:
            h += w.h(str(item), inner, s.item) + GAP_XS
        heights.append(h + g.card_inset * 2)
    card_h = max(heights) if heights else 0.0

    for i, c in enumerate(cards):
        x = g.cols_3[i]
        w.card(x, y, g.col_3_w, card_h)
        cy = y + g.card_inset
        cy += w.mono(x + g.card_inset, cy, inner, str(c.get("label", "")).upper(),
                     colour=p.accent if i == 0 else p.caption)
        cy += GAP_S
        cy += w.text(x + g.card_inset, cy, inner, str(c.get("value", "")),
                     size=s.figure_lg, bold=True, max_lines=1)
        cy += GAP_M
        cy += w.rule(x + g.card_inset, cy, inner, p.primary)
        cy += GAP_S
        for item in list(c.get("items") or [])[:5]:
            cy += w.text(x + g.card_inset, cy, inner, str(item),
                         size=s.item, colour=p.body, max_lines=2) + GAP_XS
    y += card_h
    _fits(w, "tier_cards", y)


def _metric_cards(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d)

    cards = list(d.get("cards") or [])[:4]
    inset = 0.26
    inner = g.col_4_w - inset * 2
    heights = [
        w.h_mono()
        + GAP_S + w.h(str(c.get("value", "")), inner, s.figure, bold=True)
        + GAP_S + w.h(str(c.get("caption", "")), inner, s.caption)
        + inset * 2
        for c in cards
    ]
    card_h = max(heights) if heights else 0.0

    for i, c in enumerate(cards):
        x = g.cols_4[i]
        w.card(x, y, g.col_4_w, card_h)
        cy = y + inset
        cy += w.mono(x + inset, cy, inner, str(c.get("label", "")).upper(),
                     colour=p.accent if i == 0 else p.caption)
        cy += GAP_S
        cy += w.text(x + inset, cy, inner, str(c.get("value", "")),
                     size=s.figure, bold=True, max_lines=1)
        cy += GAP_S
        w.text(x + inset, cy, inner, str(c.get("caption", "")),
               size=s.caption, colour=p.caption, max_lines=3)
    y += card_h

    chips = [str(c) for c in (d.get("chips") or [])][:4]
    if chips:
        y += GAP_L
        y += w.rule(g.margin_l, y, g.rule_w, p.primary)
        y += GAP_M
        if d.get("chips_label"):
            y += w.mono(g.margin_l, y, 3.0, str(d["chips_label"]).upper(), colour=p.caption)
            y += GAP_S
        gap = 0.16
        chip_w = (g.rule_w - (len(chips) - 1) * gap) / len(chips)
        pad = 0.18
        chip_h = max(w.h(c, chip_w - pad * 2, s.detail) for c in chips) + pad * 2
        x = g.margin_l
        for chip in chips:
            w.rect(x, y, chip_w, chip_h, p.card_alt)
            w.text(x + pad, y + pad, chip_w - pad * 2, chip, size=s.detail, max_lines=2)
            x += chip_w + gap
        y += chip_h
    _fits(w, "metric_cards", y)


def _media_feature(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d)

    media_w = 4.81
    right_l = g.margin_l + media_w + 0.40
    right_w = g.rule_w - media_w - 0.40
    ry = y
    if d.get("caption"):
        ry += w.text(right_l, ry, right_w, str(d["caption"]),
                     size=s.body_sm, colour=p.body, max_lines=4) + GAP_M

    cards = list(d.get("cards") or [])[:4]
    inset = 0.26
    gap = 0.26
    cell_w = (right_w - gap) / 2
    inner = cell_w - inset * 2
    heights = [
        w.h_mono() + GAP_S
        + w.h(str(c.get("value", "")), inner, s.name, bold=True)
        + inset * 2
        for c in cards
    ]
    cell_h = max(heights) if heights else 0.0
    for i, c in enumerate(cards):
        x = right_l + (i % 2) * (cell_w + gap)
        cy = ry + (i // 2) * (cell_h + gap)
        w.card(x, cy, cell_w, cell_h)
        ty = cy + inset
        ty += w.mono(x + inset, ty, inner, str(c.get("label", "")).upper(),
                     colour=p.accent if i == 0 else p.caption)
        ty += GAP_S
        w.text(x + inset, ty, inner, str(c.get("value", "")),
               size=s.name, bold=True, max_lines=2)
    rows = (len(cards) + 1) // 2
    ry += rows * (cell_h + gap) - (gap if rows else 0)

    media_h = max(2.40, min(CONTENT_BOTTOM, ry) - y)
    w.rect(g.margin_l, y, media_w, media_h, p.panel)
    _maybe_image(w, d, g.margin_l, y, media_w, media_h,
                 label=str(d.get("media_label", "IMAGE")))
    _fits(w, "media_feature", max(ry, y + media_h))


def _split_detail(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d)

    panels = list(d.get("panels") or [])[:2]
    gap = 0.48
    panel_w = (g.rule_w - gap) / 2 if len(panels) > 1 else g.rule_w
    inset = 0.30
    inner = panel_w - inset * 2
    heights = []
    for panel in panels:
        h = w.h_mono() + GAP_M
        for item in list(panel.get("items") or [])[:6]:
            h += w.h(str(item), inner, s.item) + GAP_S
        heights.append(h + inset * 2)
    panel_h = max(heights) if heights else 0.0

    for i, panel in enumerate(panels):
        x = g.margin_l + i * (panel_w + gap)
        w.card(x, y, panel_w, panel_h)
        cy = y + inset
        cy += w.mono(x + inset, cy, inner, str(panel.get("label", "")).upper(),
                     colour=p.accent if i == 0 else p.caption)
        cy += GAP_M
        for item in list(panel.get("items") or [])[:6]:
            cy += w.text(x + inset, cy, inner, str(item),
                         size=s.item, colour=p.body, max_lines=3) + GAP_S
    y += panel_h
    _fits(w, "split_detail", y)


def _comparison_grid(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d, size=s.headline_sm)

    columns = [str(c) for c in (d.get("columns") or [])][:5]
    rows = list(d.get("rows") or [])[:6]
    if not columns:
        return

    label_w = 3.20
    cell_w = (g.rule_w - label_w) / max(1, len(columns) - 1)
    header_h = 0.0
    for i, col in enumerate(columns):
        x = g.margin_l if i == 0 else g.margin_l + label_w + (i - 1) * cell_w
        width = (label_w if i == 0 else cell_w) - 0.16
        header_h = max(header_h, w.mono(x, y, width, col.upper(),
                                        colour=p.accent if i == 1 else p.caption))
    y += header_h + GAP_S

    for row in rows:
        y += w.rule(g.margin_l, y, g.rule_w, p.dim)
        y += GAP_S
        tallest = w.text(g.margin_l, y, label_w - 0.16, str(row.get("label", "")),
                         size=s.item, colour=p.primary, max_lines=2)
        for i, cell in enumerate([str(c) for c in (row.get("cells") or [])][:len(columns) - 1]):
            x = g.margin_l + label_w + i * cell_w
            tallest = max(tallest, w.text(x, y, cell_w - 0.16, cell, size=s.item,
                                          colour=p.accent if i == 0 else p.body,
                                          max_lines=2))
        y += tallest + GAP_S
    _fits(w, "comparison_grid", y)


def _milestone_track(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = _header(w, d)

    stages = list(d.get("stages") or [])[:4]
    if not stages:
        return
    gap = 0.30
    width = (g.rule_w - (len(stages) - 1) * gap) / len(stages)
    inset = 0.28
    inner = width - inset * 2
    card_h = max(
        w.h_mono() + GAP_S
        + w.h(str(st.get("status", "")), inner, s.name, bold=True) + GAP_S
        + w.h(str(st.get("detail", "")), inner, s.detail_sm)
        + inset * 2
        for st in stages
    )

    for i, st in enumerate(stages):
        x = g.margin_l + i * (width + gap)
        done = bool(st.get("done"))
        w.card(x, y, width, card_h)
        w.rule(x, y, width, p.accent if done else p.dim)
        cy = y + inset
        cy += w.mono(x + inset, cy, inner, str(st.get("label", "")).upper(),
                     colour=p.accent if done else p.caption)
        cy += GAP_S
        cy += w.text(x + inset, cy, inner, str(st.get("status", "")),
                     size=s.name, bold=True, max_lines=2)
        cy += GAP_S
        w.text(x + inset, cy, inner, str(st.get("detail", "")),
               size=s.detail_sm, colour=p.caption, max_lines=5)
    y += card_h
    _fits(w, "milestone_track", y)


def _hero_number(w: _SlideWriter, d: dict[str, Any]) -> None:
    g, p, s = w.t.grid, w.t.palette, w.t.sizes
    y = 2.03
    eb = w.eyebrow(y, str(d.get("eyebrow", "")))
    if eb:
        y += eb + GAP_M
    y += w.text(g.margin_l, y, g.rule_w, str(d.get("headline", "")),
                size=s.hero, bold=True, spacing=0.95, max_lines=1)
    y += GAP_M
    if d.get("line"):
        y += w.text(g.margin_l, y, g.rule_w, str(d["line"]),
                    size=s.figure, bold=True, colour=p.accent, max_lines=2)
        y += GAP_M
    if d.get("caption"):
        y += w.text(g.margin_l, y, 6.0, str(d["caption"]),
                    size=s.detail, colour=p.caption, max_lines=3)
    _fits(w, "hero_number", y)


def _maybe_image(w: _SlideWriter, d, left, top, width, height, *, label: str | None = None):
    """Place an image if one resolves, otherwise leave the labelled panel."""
    src = d.get("image")
    if src:
        path = Path(str(src))
        if path.is_file():
            try:
                w.s.shapes.add_picture(str(path), _in(left), _in(top),
                                       _in(width), _in(height))
                return
            except Exception:  # noqa: BLE001 — a bad image must not kill the deck
                logger.exception("Could not place image %s", src)
    if label:
        w.text(left, top + height / 2 - 0.11, width, label,
               font=w.t.fonts.mono, size=w.t.sizes.eyebrow,
               colour=w.t.palette.dim, align=PP_ALIGN.CENTER)


_RENDERERS = {
    "cover": _cover,
    "section_break": _section_break,
    "statement_media": _statement_media,
    "bullets": _bullets,
    "points_with_metrics": _points_with_metrics,
    "numbered_cards": _numbered_cards,
    "profile_cards": _profile_cards,
    "tier_cards": _tier_cards,
    "metric_cards": _metric_cards,
    "media_feature": _media_feature,
    "split_detail": _split_detail,
    "comparison_grid": _comparison_grid,
    "milestone_track": _milestone_track,
    "hero_number": _hero_number,
}


def build_pptx(spec: dict[str, Any], out_path: Path, theme: DeckTheme | None = None) -> Path:
    """Render a deck spec to a .pptx file. Raises DeckBuildError on a bad spec."""
    theme = theme or DEFAULT_THEME
    slides = spec.get("slides") or []
    if not slides:
        raise DeckBuildError("The deck needs at least one slide.")

    unknown = {s.get("archetype") for s in slides} - set(_RENDERERS)
    if unknown:
        raise DeckBuildError(
            f"Unknown archetype(s): {', '.join(sorted(str(u) for u in unknown))}. "
            f"Available: {', '.join(sorted(_RENDERERS))}."
        )

    # Required, not defaulted: how a deck is marked is the user's call, and
    # guessing it wrong is the kind of mistake that only shows up outside.
    marks = theme.classifications
    key = str(spec.get("classification", "")).strip().lower().replace(" ", "_").replace("-", "_")
    if not key:
        raise DeckBuildError(
            "classification is required. Ask the user which applies: "
            f"{', '.join(marks)}."
        )
    if key not in marks:
        raise DeckBuildError(
            f'"{key}" is not a classification. Use one of: {", ".join(marks)}.'
        )
    mark = marks[key]

    prs = Presentation()
    prs.slide_width = theme.canvas_w
    prs.slide_height = theme.canvas_h
    blank = prs.slide_layouts[6]

    for i, data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(theme.palette.background)
        w = _SlideWriter(slide, theme)
        name = str(data.get("archetype"))
        # The cover carries no page number, matching the example.
        w.chrome(None if i == 0 and name == "cover" else i + 1, mark)
        _RENDERERS[name](w, data)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def archetype_names() -> list[str]:
    return sorted(_RENDERERS)


assert set(_RENDERERS) == set(BY_NAME), "renderers and catalog must stay in step"
